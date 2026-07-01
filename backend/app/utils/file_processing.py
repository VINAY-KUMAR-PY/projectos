from __future__ import annotations

import csv
import json
from pathlib import Path
import zipfile


CODE_EXTENSIONS = {
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".java",
    ".go",
    ".rs",
    ".php",
    ".rb",
    ".cs",
    ".cpp",
    ".c",
    ".html",
    ".css",
    ".md",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
}

TECH_STACK_MARKERS = {
    "package.json": "Node.js",
    "next.config.js": "Next.js",
    "next.config.mjs": "Next.js",
    "requirements.txt": "Python",
    "pyproject.toml": "Python",
    "manage.py": "Django",
    "main.py": "Python",
    "app.py": "Python",
    "pom.xml": "Java",
    "Cargo.toml": "Rust",
    "go.mod": "Go",
    "Dockerfile": "Docker",
}


def validate_upload_name(filename: str, allowed_extensions: set[str] | None = None) -> str:
    """Return a safe upload filename after validating its extension."""
    safe_name = Path(filename or "upload.bin").name
    allowed = allowed_extensions or {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xls",
        ".csv",
        ".png",
        ".jpg",
        ".jpeg",
        ".webp",
        ".gif",
        ".zip",
        ".mp4",
        ".mov",
        ".mp3",
        ".wav",
        ".txt",
        ".md",
        ".py",
        ".js",
        ".ts",
        ".json",
        ".yaml",
        ".yml",
    }
    suffix = Path(safe_name).suffix.lower()
    if suffix not in allowed:
        raise ValueError(f"File type {suffix or '<none>'} is not allowed")
    return safe_name


def parse_uploaded_file(file_path: str, file_type: str | None = None) -> dict:
    """Parse supported upload types into text plus useful metadata."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    file_type = file_type or ""
    metadata: dict = {
        "file_name": path.name,
        "extension": suffix,
        "content_type": file_type,
    }
    text = ""
    parser = "fallback"

    if file_type.startswith("text/") or suffix in CODE_EXTENSIONS:
        text = path.read_text(encoding="utf-8", errors="ignore")
        parser = "text"
        metadata["line_count"] = len(text.splitlines())
        metadata["tech_stack"] = detect_tech_stack([path.name])

    elif suffix == ".pdf" or file_type == "application/pdf":
        text, metadata, parser = _parse_pdf(path, metadata)

    elif suffix == ".docx":
        text, metadata, parser = _parse_docx(path, metadata)

    elif suffix == ".pptx":
        text, metadata, parser = _parse_pptx(path, metadata)

    elif suffix in {".xlsx", ".xls", ".csv"}:
        text, metadata, parser = _parse_table(path, metadata)

    elif file_type.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        text, metadata, parser = _parse_image(path, metadata)

    elif suffix == ".zip" or file_type in {"application/zip", "application/x-zip-compressed"}:
        text, metadata, parser = _parse_zip(path, metadata)

    else:
        text = "[File type not supported for text extraction]"

    metadata["parser"] = parser
    metadata["text_length"] = len(text)
    return {
        "text": text,
        "summary": summarize_text(text),
        "metadata": metadata,
    }


def extract_text_from_file(file_path: str, file_type: str | None = None) -> str:
    return parse_uploaded_file(file_path, file_type)["text"]


def extract_video_transcript(file_path: str) -> str:
    try:
        import whisper

        model = whisper.load_model("base")
        result = model.transcribe(file_path)
        return result.get("text", "")
    except Exception as exc:
        return f"[Transcript extraction unavailable: {exc}]"


def extract_key_frames(file_path: str, every_seconds: int = 10, max_frames: int = 5) -> list[dict]:
    """Extract a small set of frame image paths if OpenCV is installed."""
    try:
        import cv2
    except Exception as exc:
        return [{"status": "unavailable", "detail": f"OpenCV not installed: {exc}"}]

    capture = cv2.VideoCapture(file_path)
    fps = capture.get(cv2.CAP_PROP_FPS) or 30
    interval = max(1, int(fps * every_seconds))
    frames = []
    frame_index = 0
    output_dir = Path(file_path).with_suffix("")
    output_dir.mkdir(parents=True, exist_ok=True)
    while len(frames) < max_frames:
        ok, frame = capture.read()
        if not ok:
            break
        if frame_index % interval == 0:
            target = output_dir / f"frame_{frame_index}.jpg"
            cv2.imwrite(str(target), frame)
            frames.append({"frame": frame_index, "path": str(target)})
        frame_index += 1
    capture.release()
    return frames


def detect_tech_stack(paths: list[str]) -> list[str]:
    detected = set()
    for raw_path in paths:
        name = Path(raw_path).name
        suffix = Path(raw_path).suffix.lower()
        if name in TECH_STACK_MARKERS:
            detected.add(TECH_STACK_MARKERS[name])
        if suffix in {".ts", ".tsx"}:
            detected.add("TypeScript")
        if suffix in {".js", ".jsx"}:
            detected.add("JavaScript")
        if suffix == ".py":
            detected.add("Python")
        if suffix in {".html", ".css"}:
            detected.add("Frontend")
    return sorted(detected)


def summarize_text(text: str, max_chars: int = 1000) -> str:
    compact = " ".join(text.split())
    return compact[:max_chars] if compact else "[No extractable text]"


def _parse_pdf(path: Path, metadata: dict) -> tuple[str, dict, str]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        metadata["page_count"] = len(reader.pages)
        return "\n".join(pages), metadata, "pypdf"
    except Exception as first_exc:
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
                metadata["page_count"] = len(pdf.pages)
            return "\n".join(pages), metadata, "pdfplumber"
        except Exception as second_exc:
            return (
                f"[PDF extraction unavailable: {first_exc}; {second_exc}]",
                metadata,
                "pdf-unavailable",
            )


def _parse_docx(path: Path, metadata: dict) -> tuple[str, dict, str]:
    try:
        from docx import Document

        document = Document(str(path))
        headings = [
            paragraph.text
            for paragraph in document.paragraphs
            if paragraph.style and paragraph.style.name.lower().startswith("heading")
        ]
        metadata["headings"] = headings
        return "\n".join(paragraph.text for paragraph in document.paragraphs), metadata, "python-docx"
    except Exception as exc:
        return f"[DOCX extraction unavailable: {exc}]", metadata, "docx-unavailable"


def _parse_pptx(path: Path, metadata: dict) -> tuple[str, dict, str]:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            slides.append({"slide": index, "text": "\n".join(parts)})
        metadata["slide_count"] = len(slides)
        metadata["slides"] = slides
        return "\n\n".join(f"Slide {item['slide']}:\n{item['text']}" for item in slides), metadata, "python-pptx"
    except Exception as exc:
        return f"[PPTX extraction unavailable: {exc}]", metadata, "pptx-unavailable"


def _parse_table(path: Path, metadata: dict) -> tuple[str, dict, str]:
    suffix = path.suffix.lower()
    try:
        import pandas as pd

        frame = pd.read_csv(path) if suffix == ".csv" else pd.read_excel(path)
        metadata["columns"] = list(frame.columns)
        metadata["row_count"] = int(len(frame))
        metadata["preview"] = frame.head(10).to_dict(orient="records")
        return frame.head(50).to_string(index=False), metadata, "pandas"
    except Exception:
        if suffix == ".csv":
            with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
                rows = list(csv.reader(handle))
            metadata["row_count"] = max(0, len(rows) - 1)
            metadata["columns"] = rows[0] if rows else []
            return "\n".join(",".join(row) for row in rows[:50]), metadata, "csv"
        return "[Spreadsheet extraction unavailable: install pandas/openpyxl]", metadata, "table-unavailable"


def _parse_image(path: Path, metadata: dict) -> tuple[str, dict, str]:
    try:
        import pytesseract
        from PIL import Image

        ocr_text = pytesseract.image_to_string(Image.open(path))
        metadata["vision_description"] = summarize_text(ocr_text, 300)
        return ocr_text or "[No OCR text detected]", metadata, "pytesseract"
    except Exception as exc:
        metadata["vision_description"] = "Image uploaded; OCR optional dependency unavailable."
        return f"[Image OCR unavailable: {exc}]", metadata, "ocr-unavailable"


def _parse_zip(path: Path, metadata: dict) -> tuple[str, dict, str]:
    extracted = []
    file_tree = []
    with zipfile.ZipFile(path, "r") as archive:
        for name in archive.namelist():
            if name.endswith("/"):
                continue
            file_tree.append(name)
            if Path(name).suffix.lower() in CODE_EXTENSIONS:
                with archive.open(name) as item:
                    content = item.read().decode("utf-8", errors="ignore")
                    extracted.append(f"=== {name} ===\n{content[:5000]}")
    metadata["file_tree"] = file_tree[:500]
    metadata["file_count"] = len(file_tree)
    metadata["tech_stack"] = detect_tech_stack(file_tree)
    return "\n\n".join(extracted) or json.dumps(metadata["file_tree"][:100]), metadata, "zip"
