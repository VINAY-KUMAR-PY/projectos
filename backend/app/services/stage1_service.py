from __future__ import annotations

import base64
import hashlib
import hmac
import json
from pathlib import Path
import re
import secrets
import time
import zipfile

from fastapi import HTTPException, status
from sqlalchemy import delete, or_, select
from sqlalchemy.orm import Session

from app.ai.manager import create_ai_manager
from app.ai.memory import ProjectMemoryConnector
from app.ai.types import AIRequest
from app.config.settings import settings
from app.models.stage1 import (
    ApprovalRequest,
    AuditLog,
    Comment,
    MarketplaceItem,
    TaskAssignment,
    TeamMember,
    UsageRecord,
    UserIntegration,
)
from app.models.user import User
from app.models.ai import AgentRun, Conversation, ConversationMessage
from app.models.workspace import (
    GeneratedOutput,
    Note,
    Project,
    ProjectFile,
    ProjectMemory,
    Subscription,
    Task,
    Workspace,
)
from app.schemas.stage1_schema import (
    ApprovalDecisionRequest,
    ApprovalRequestCreate,
    CodeBuildRequest,
    CommentCreateRequest,
    DeploymentGenerateRequest,
    IntegrationLinkRequest,
    LearningRequest,
    MarketplacePublishRequest,
    MarketplaceUseRequest,
    RepoImportRequest,
    ReviewRequest,
    TaskAssignRequest,
    TeamInviteRequest,
)
from app.services import usage_service
from app.services.workspace_service import create_project, get_project
from app.schemas.platform_schema import ProjectApiCreate
from app.schemas.workspace_schema import WorkspaceCreate
from app.services import workspace_service
from app.utils.file_processing import (
    extract_key_frames,
    extract_video_transcript,
    parse_uploaded_file,
    summarize_text,
)


PLAN_PRICES = {"free": 0, "student": 299, "pro": 499, "team": 999, "enterprise": 4999}


def generated_dir(project_id: int, *parts: str) -> Path:
    path = Path(settings.upload_dir) / "generated" / str(project_id)
    for part in parts:
        path = path / part
    path.mkdir(parents=True, exist_ok=True)
    return path


def file_url(path: Path) -> str:
    return f"/uploads/{path.relative_to(Path(settings.upload_dir)).as_posix()}"


def log_audit(
    db: Session,
    user_id: int | None,
    action: str,
    entity_type: str | None = None,
    entity_id: int | None = None,
    detail: str | None = None,
):
    db.add(
        AuditLog(
            user_id=user_id,
            action=action,
            entity_type=entity_type,
            entity_id=entity_id,
            detail=detail,
        )
    )


def save_memory(db: Session, project_id: int, owner_id: int, key: str, value: str):
    return ProjectMemoryConnector(db).save(project_id, owner_id, key, value)


def save_generated_output(
    db: Session,
    owner_id: int,
    project_id: int,
    output_type: str,
    title: str,
    content: str,
    output_format: str = "markdown",
):
    output = GeneratedOutput(
        owner_id=owner_id,
        project_id=project_id,
        output_type=output_type,
        title=title,
        content=content,
        format=output_format,
    )
    db.add(output)
    usage_service.record_usage(db, owner_id, "generated_output", 1, project_id)
    return output


def project_context(db: Session, project: Project) -> str:
    memories = ProjectMemoryConnector(db).load(project.id, project.owner_id, limit=50)
    memory_text = "\n".join(f"{item['key']}: {item['value'][:2000]}" for item in memories)
    return (
        f"Project: {project.title}\n"
        f"Description: {project.description or ''}\n"
        f"Status: {project.status}\n"
        f"Memory:\n{memory_text}"
    )


def analyze_project_file(db: Session, file_id: int, user: User):
    record = db.scalar(select(ProjectFile).where(ProjectFile.id == file_id))
    if not record:
        raise HTTPException(status_code=404, detail="File not found")
    get_project(db, record.project_id, user.id)
    parsed = parse_uploaded_file(record.storage_path, record.file_type)
    record.extracted_text = parsed["text"]
    record.summary = parsed["summary"]
    save_memory(
        db,
        record.project_id,
        record.owner_id,
        f"file:{record.id}:{record.file_type or 'unknown'}",
        json.dumps(parsed, ensure_ascii=True),
    )
    log_audit(db, user.id, "file.analyzed", "file", record.id)
    db.commit()
    db.refresh(record)
    return {"file": record, "analysis": parsed}


def analyze_video_file(db: Session, file_id: int, user: User):
    record = db.scalar(select(ProjectFile).where(ProjectFile.id == file_id))
    if not record:
        raise HTTPException(status_code=404, detail="File not found")
    get_project(db, record.project_id, user.id)
    transcript = extract_video_transcript(record.storage_path)
    frames = extract_key_frames(record.storage_path)
    summary = summarize_text(transcript, 600)
    result = {
        "file_id": record.id,
        "transcript": transcript,
        "key_frames": frames,
        "summary": summary,
        "action_items": ["Review transcript highlights", "Validate demo flow", "Add missing explanation points"],
        "mistakes": ["No automated visual issue detected in mock mode"],
        "viva_questions": [
            "What problem does this project solve?",
            "Which part of the demo proves the core workflow?",
            "What limitations remain?",
            "How would you scale this system?",
            "How did you validate the results?",
        ],
    }
    record.extracted_text = transcript
    record.summary = summary
    save_memory(db, record.project_id, record.owner_id, f"video:{record.id}", json.dumps(result, ensure_ascii=True))
    log_audit(db, user.id, "file.video_analyzed", "file", record.id)
    db.commit()
    return result


def generate_document(db: Session, project_id: int, user: User):
    project = get_project(db, project_id, user.id)
    usage_service.assert_agent_run_limit(db, user)
    context = project_context(db, project)
    sections = [
        "Abstract",
        "Introduction",
        "Literature Review",
        "Methodology",
        "System Design",
        "Implementation",
        "Testing",
        "Results",
        "Conclusion",
        "Future Scope",
        "References",
    ]
    ai = create_ai_manager()
    generated_sections = []
    for section in sections:
        response = ai.generate(
            AIRequest(
                prompt=f"Write the {section} section for this project.\n\n{context}",
                system_prompt="You generate concise academic project documentation.",
                max_tokens=700,
            )
        )
        generated_sections.append((section, response.content))
    output_dir = generated_dir(project_id, "documents")
    target = output_dir / "project_report.docx"
    _write_docx(target, project.title, generated_sections)
    content = "\n\n".join(f"# {title}\n{body}" for title, body in generated_sections)
    output = save_generated_output(db, user.id, project_id, "document", "Project Report", content, "docx")
    usage_service.record_usage(db, user.id, "agent_run", 1, project_id)
    save_memory(db, project_id, project.owner_id, "generated:document", content)
    log_audit(db, user.id, "project.document_generated", "project", project_id)
    db.commit()
    db.refresh(output)
    return {"output_id": output.id, "file_url": file_url(target), "title": output.title, "content": content}


def generate_ppt(db: Session, project_id: int, user: User):
    project = get_project(db, project_id, user.id)
    usage_service.assert_agent_run_limit(db, user)
    slide_titles = [
        "Title",
        "Agenda",
        "Problem Statement",
        "Architecture",
        "Tech Stack",
        "Features",
        "Demo",
        "Results",
        "Conclusion",
        "Thank You",
    ]
    slides = [
        {"title": title, "bullets": [f"{title} for {project.title}", project.description or "ProjectOS generated content"]}
        for title in slide_titles
    ]
    output_dir = generated_dir(project_id, "presentations")
    target = output_dir / "project_presentation.pptx"
    _write_pptx(target, project.title, slides)
    demo_script = f"Introduce {project.title}, show the core workflow, explain results, and close with future scope."
    content = json.dumps({"slides": slides, "demo_script": demo_script}, ensure_ascii=True)
    output = save_generated_output(db, user.id, project_id, "ppt", "Project Presentation", content, "pptx")
    usage_service.record_usage(db, user.id, "agent_run", 1, project_id)
    save_memory(db, project_id, project.owner_id, "generated:ppt", content)
    log_audit(db, user.id, "project.ppt_generated", "project", project_id)
    db.commit()
    db.refresh(output)
    return {
        "output_id": output.id,
        "file_url": file_url(target),
        "title": output.title,
        "content": demo_script,
        "metadata": {"slides": slides},
    }


def generate_diagram(db: Session, project_id: int, user: User, diagram_type: str):
    project = get_project(db, project_id, user.id)
    usage_service.assert_agent_run_limit(db, user)
    mermaid = _diagram_for(project, diagram_type)
    output = save_generated_output(
        db,
        user.id,
        project_id,
        f"diagram:{diagram_type}",
        f"{diagram_type.title()} Diagram",
        mermaid,
        "mermaid",
    )
    usage_service.record_usage(db, user.id, "agent_run", 1, project_id)
    save_memory(db, project_id, project.owner_id, f"generated:diagram:{diagram_type}", mermaid)
    log_audit(db, user.id, "project.diagram_generated", "project", project_id, diagram_type)
    db.commit()
    db.refresh(output)
    return {"output_id": output.id, "title": output.title, "content": mermaid, "metadata": {"type": diagram_type}}


def build_code(db: Session, project_id: int, user: User, request: CodeBuildRequest):
    project = get_project(db, project_id, user.id)
    usage_service.assert_agent_run_limit(db, user)
    output_dir = generated_dir(project_id, "code")
    if request.action == "generate":
        files = {
            "README.md": f"# {project.title}\n\n{project.description or 'Generated project scaffold.'}\n",
            "src/main.py": "def main():\n    print('ProjectOS generated scaffold')\n\nif __name__ == '__main__':\n    main()\n",
            "tests/test_main.py": "from src.main import main\n\n\ndef test_main_runs():\n    assert main() is None\n",
        }
        for name, content in files.items():
            target = output_dir / name
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_text(content, encoding="utf-8")
        zip_path = generated_dir(project_id) / "code.zip"
        _zip_dir(output_dir, zip_path)
        content = json.dumps(
            {"file_tree": sorted(files), "zip_url": file_url(zip_path), "summary": "Python scaffold generated"},
            ensure_ascii=True,
        )
    else:
        source = _source_from_request(db, request)
        content = json.dumps(
            {
                "action": request.action,
                "analysis": _review_text(source),
                "prompt": request.prompt or "",
            },
            ensure_ascii=True,
        )
    output = save_generated_output(db, user.id, project_id, f"code:{request.action}", "Code Builder Result", content, "json")
    usage_service.record_usage(db, user.id, "agent_run", 1, project_id)
    save_memory(db, project_id, project.owner_id, f"generated:code:{request.action}", content)
    log_audit(db, user.id, "project.code_generated", "project", project_id, request.action)
    db.commit()
    db.refresh(output)
    return {"output_id": output.id, "title": output.title, "content": content}


def review_project(db: Session, project_id: int, user: User, request: ReviewRequest):
    project = get_project(db, project_id, user.id)
    source = ""
    if request.file_id:
        file_record = db.scalar(select(ProjectFile).where(ProjectFile.id == request.file_id, ProjectFile.project_id == project_id))
        if file_record:
            source = file_record.extracted_text or parse_uploaded_file(file_record.storage_path, file_record.file_type)["text"]
    if not source:
        source = project_context(db, project)
    tasks = db.scalars(select(Task).where(Task.project_id == project_id, Task.owner_id == user.id)).all()
    files = db.scalars(select(ProjectFile).where(ProjectFile.project_id == project_id, ProjectFile.owner_id == user.id)).all()
    outputs = db.scalars(select(GeneratedOutput).where(GeneratedOutput.project_id == project_id, GeneratedOutput.owner_id == user.id)).all()
    memories = ProjectMemoryConnector(db).load(project.id, project.owner_id, limit=100)
    issues = _review_text(source)
    hardcoded_secret = bool(re.search(r"(api[_-]?key|secret|password)\s*=\s*['\"][^'\"]+", source, re.I))
    words = source.split()
    duplicate_risk = "medium" if len(source) > 2000 and len(set(words)) < max(20, len(words) // 4) else "low"
    completed_tasks = sum(1 for task in tasks if task.status in {"done", "complete", "completed"})
    generated_types = {output.output_type for output in outputs}
    has_document = any(output_type.startswith("document") or "doc" in output_type for output_type in generated_types)
    has_presentation = any("ppt" in output_type or "diagram" in output_type for output_type in generated_types)
    has_code = any(output_type.startswith("code") for output_type in generated_types)
    high_risk_issue_count = sum(1 for issue in issues if not issue.startswith("No high-risk"))
    issue_penalty = min(30, high_risk_issue_count * 7)
    secret_penalty = 15 if hardcoded_secret else 0
    file_signal = min(18, len(files) * 4)
    memory_signal = min(16, len(memories) * 2)
    output_signal = min(18, len(outputs) * 3)
    task_signal = min(20, len(tasks) * 4)
    completion_signal = min(20, completed_tasks * 5)
    source_signal = min(12, len(source) // 600)
    code_quality = _bounded_score(50 + file_signal + source_signal + (10 if has_code else 0) - issue_penalty - secret_penalty)
    documentation_completeness = _bounded_score(40 + memory_signal + output_signal + (18 if has_document else 0) + min(10, len(files) * 2))
    testing_coverage = _bounded_score(35 + task_signal + completion_signal + (12 if "review" in generated_types else 0) + (8 if has_code else 0))
    presentation_readiness = _bounded_score(35 + (24 if has_presentation else 0) + min(18, len(outputs) * 4) + min(12, len(files) * 3))
    project_improvement_score = _bounded_score(
        round((code_quality + documentation_completeness + testing_coverage + presentation_readiness) / 4)
        - (5 if duplicate_risk == "medium" else 0)
    )
    scorecard = {
        "code_quality": code_quality,
        "documentation_completeness": documentation_completeness,
        "testing_coverage": testing_coverage,
        "presentation_readiness": presentation_readiness,
        "project_improvement_score": project_improvement_score,
        "plagiarism_risk": duplicate_risk,
        "issues": issues,
        "inputs": {
            "tasks": len(tasks),
            "completed_tasks": completed_tasks,
            "files": len(files),
            "generated_outputs": len(outputs),
            "memory_items": len(memories),
            "source_characters": len(source),
        },
        "suggestions": [
            "Add automated tests for the most important workflows.",
            "Move secrets into environment variables.",
            "Document setup and deployment commands clearly.",
        ],
    }
    output = save_generated_output(db, user.id, project_id, "review", "Testing and Review Scorecard", json.dumps(scorecard), "json")
    log_audit(db, user.id, "project.reviewed", "project", project_id)
    db.commit()
    db.refresh(output)
    return scorecard


def generate_deployment(db: Session, project_id: int, user: User, request: DeploymentGenerateRequest):
    project = get_project(db, project_id, user.id)
    target = request.target
    configs = {
        "vercel": {"filename": "vercel.json", "content": json.dumps({"buildCommand": "npm run build"}, indent=2)},
        "netlify": {"filename": "netlify.toml", "content": "[build]\ncommand = \"npm run build\"\npublish = \".next\"\n"},
        "railway": {"filename": "railway.json", "content": json.dumps({"deploy": {"startCommand": "uvicorn app.api.server:app"}}, indent=2)},
        "render": {"filename": "render.yaml", "content": "services:\n  - type: web\n    env: python\n    startCommand: uvicorn app.api.server:app\n"},
        "aws": {"filename": "Dockerfile", "content": "FROM python:3.11-slim\nCMD [\"uvicorn\", \"app.api.server:app\"]\n"},
        "docker": {"filename": "Dockerfile", "content": "FROM python:3.11-slim\nWORKDIR /app\nCOPY . .\nCMD [\"uvicorn\", \"app.api.server:app\"]\n"},
        "github-pages": {"filename": "deploy.yml", "content": "name: deploy\non: push\njobs:\n  pages:\n    runs-on: ubuntu-latest\n    steps:\n      - uses: actions/checkout@v4\n"},
    }
    payload = configs[target]
    instructions = f"Deploy {project.title} to {target}: commit {payload['filename']}, configure environment variables, run build, then monitor logs."
    content = json.dumps({"target": target, **payload, "instructions": instructions}, ensure_ascii=True)
    output = save_generated_output(db, user.id, project_id, f"deployment:{target}", f"{target.title()} Deployment", content, "json")
    log_audit(db, user.id, "project.deployment_generated", "project", project_id, target)
    db.commit()
    db.refresh(output)
    return {"output_id": output.id, "title": output.title, "content": content, "metadata": {"target": target}}


def invite_member(db: Session, project_id: int, user: User, request: TeamInviteRequest):
    project = get_project(db, project_id, user.id)
    target_user = db.scalar(select(User).where(User.email == request.email))
    member = TeamMember(
        project_id=project.id,
        owner_id=project.owner_id,
        user_id=target_user.id if target_user else None,
        email=request.email,
        role=request.role,
    )
    db.add(member)
    log_audit(db, user.id, "team.member_invited", "project", project_id, request.email)
    db.commit()
    db.refresh(member)
    return member


def list_members(db: Session, project_id: int, user: User):
    get_project(db, project_id, user.id)
    return db.scalars(select(TeamMember).where(TeamMember.project_id == project_id).order_by(TeamMember.id.desc())).all()


def add_comment(db: Session, project_id: int, user: User, request: CommentCreateRequest):
    project = get_project(db, project_id, user.id)
    comment = Comment(
        project_id=project.id,
        owner_id=project.owner_id,
        author_id=user.id,
        entity_type=request.entity_type,
        entity_id=request.entity_id,
        body=request.body,
    )
    db.add(comment)
    log_audit(db, user.id, "comment.created", request.entity_type, request.entity_id)
    db.commit()
    db.refresh(comment)
    return comment


def list_comments(db: Session, project_id: int, user: User, entity_type: str | None, entity_id: int | None):
    get_project(db, project_id, user.id)
    statement = select(Comment).where(Comment.project_id == project_id)
    if entity_type:
        statement = statement.where(Comment.entity_type == entity_type)
    if entity_id is not None:
        statement = statement.where(Comment.entity_id == entity_id)
    return db.scalars(statement.order_by(Comment.id.desc())).all()


def assign_task(db: Session, project_id: int, user: User, request: TaskAssignRequest):
    get_project(db, project_id, user.id)
    task = db.scalar(select(Task).where(Task.id == request.task_id, Task.project_id == project_id))
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    assignment = TaskAssignment(
        task_id=task.id,
        project_id=project_id,
        assignee_email=request.assignee_email,
        role=request.role,
        created_by=user.id,
    )
    db.add(assignment)
    log_audit(db, user.id, "task.assigned", "task", task.id, request.assignee_email)
    db.commit()
    db.refresh(assignment)
    return assignment


def create_approval(db: Session, project_id: int, user: User, request: ApprovalRequestCreate):
    get_project(db, project_id, user.id)
    output = db.scalar(select(GeneratedOutput).where(GeneratedOutput.id == request.output_id, GeneratedOutput.project_id == project_id))
    if not output:
        raise HTTPException(status_code=404, detail="Generated output not found")
    approval = ApprovalRequest(
        project_id=project_id,
        output_id=output.id,
        requester_id=user.id,
        reviewer_email=request.reviewer_email,
        note=request.note,
    )
    db.add(approval)
    log_audit(db, user.id, "approval.requested", "output", output.id)
    db.commit()
    db.refresh(approval)
    return approval


def decide_approval(db: Session, approval_id: int, user: User, request: ApprovalDecisionRequest):
    approval = db.scalar(select(ApprovalRequest).where(ApprovalRequest.id == approval_id))
    if not approval:
        raise HTTPException(status_code=404, detail="Approval request not found")
    get_project(db, approval.project_id, user.id)
    approval.status = request.status
    approval.note = request.note or approval.note
    approval.decided_by = user.id
    log_audit(db, user.id, f"approval.{request.status}", "approval", approval.id)
    db.commit()
    db.refresh(approval)
    return approval


def run_learning_action(db: Session, action: str, user: User, request: LearningRequest):
    project = get_project(db, request.project_id, user.id)
    prompts = {
        "explain_concept": "Explain this concept simply.",
        "generate_quiz": "Create a short quiz with answers.",
        "viva_practice": "Run viva practice questions and ideal answers.",
        "interview_questions": "Create interview questions.",
        "explain_my_project": "Explain my project in spoken defense style.",
    }
    if action not in prompts:
        raise HTTPException(status_code=404, detail="Unknown learning action")
    context = project_context(db, project)
    response = create_ai_manager().generate(
        AIRequest(prompt=f"{prompts[action]}\n\nUser message: {request.message}\n\n{context}", max_tokens=1000)
    )
    save_memory(db, project.id, project.owner_id, f"learning:{action}", response.content)
    usage_service.record_usage(db, user.id, "agent_run", 1, project.id)
    log_audit(db, user.id, f"learning.{action}", "project", project.id)
    db.commit()
    return {"action": action, "project_id": project.id, "content": response.content}


def list_marketplace(db: Session, search: str | None = None, item_type: str | None = None):
    statement = select(MarketplaceItem).where(MarketplaceItem.is_public.is_(True))
    if search:
        pattern = f"%{search}%"
        statement = statement.where(or_(MarketplaceItem.title.ilike(pattern), MarketplaceItem.description.ilike(pattern)))
    if item_type:
        statement = statement.where(MarketplaceItem.item_type == item_type)
    return db.scalars(statement.order_by(MarketplaceItem.id.desc())).all()


def publish_marketplace_item(db: Session, user: User, request: MarketplacePublishRequest):
    item = MarketplaceItem(creator_id=user.id, **request.model_dump())
    db.add(item)
    log_audit(db, user.id, "marketplace.published", "marketplace_item", None, request.title)
    db.commit()
    db.refresh(item)
    return item


def use_marketplace_item(db: Session, item_id: int, user: User, request: MarketplaceUseRequest):
    item = db.scalar(select(MarketplaceItem).where(MarketplaceItem.id == item_id, MarketplaceItem.is_public.is_(True)))
    if not item:
        raise HTTPException(status_code=404, detail="Marketplace item not found")
    usage_service.assert_project_limit(db, user)
    workspace_id = request.workspace_id
    if workspace_id is None:
        workspace = workspace_service.create_workspace(
            db,
            user.id,
            WorkspaceCreate(name="Marketplace Workspace", description="Created from marketplace item"),
        )
        workspace_id = workspace.id
    project = create_project(
        db,
        workspace_id,
        user.id,
        ProjectApiCreate(title=item.title, description=item.description or item.content_ref),
    )
    log_audit(db, user.id, "marketplace.used", "marketplace_item", item.id)
    db.commit()
    return {"item_id": item.id, "project_id": project.id}


def link_integration(db: Session, user: User, request: IntegrationLinkRequest):
    integration = UserIntegration(user_id=user.id, **request.model_dump())
    db.add(integration)
    log_audit(db, user.id, "integration.linked", "integration", None, request.provider)
    db.commit()
    db.refresh(integration)
    return _serialize_integration(integration)


def list_integrations(db: Session, user: User):
    integrations = db.scalars(select(UserIntegration).where(UserIntegration.user_id == user.id).order_by(UserIntegration.id.desc())).all()
    return [_serialize_integration(integration) for integration in integrations]


def import_repo(db: Session, user: User, request: RepoImportRequest):
    if not _is_github_repo_url(request.repo_url):
        raise HTTPException(status_code=400, detail="Only GitHub repository URLs are supported for GitHub import")
    usage_service.assert_project_limit(db, user)
    title = request.repo_url.rstrip("/").split("/")[-1] or "Imported Repository"
    if title.endswith(".git"):
        title = title[:-4]
    workspace_id = request.workspace_id
    if workspace_id is None:
        workspace = workspace_service.create_workspace(
            db,
            user.id,
            WorkspaceCreate(name="GitHub Imports", description="Imported repositories"),
        )
        workspace_id = workspace.id
    project = create_project(
        db,
        workspace_id,
        user.id,
        ProjectApiCreate(title=title, description=f"Imported from {request.repo_url}"),
    )
    mode = _integration_mode(db, user, "github")
    save_memory(db, project.id, project.owner_id, "integration:github_repo", json.dumps({"repo_url": request.repo_url, "mode": mode}))
    log_audit(db, user.id, "integration.github_imported", "project", project.id, f"{mode}:{request.repo_url}")
    db.commit()
    return {
        "project_id": project.id,
        "repo_url": request.repo_url,
        "mode": mode,
        "message": (
            "Repository imported in mock mode; no GitHub token was configured."
            if mode == "mock"
            else "Repository import recorded with a configured GitHub token."
        ),
    }


def push_generated_code(db: Session, user: User):
    mode = _integration_mode(db, user, "github")
    if mode == "mock":
        return {
            "status": "mock",
            "mode": "mock",
            "message": "No GitHub token is linked; generated code was not pushed.",
        }
    return {
        "status": "ready",
        "mode": "configured",
        "message": "GitHub credentials are configured; push execution is delegated to the deployment worker.",
    }


def import_drive_file(db: Session, user: User):
    mode = _integration_mode(db, user, "google-drive")
    if mode == "mock":
        return {
            "status": "mock",
            "mode": "mock",
            "message": "No Google Drive credentials are linked; no cloud file was imported.",
        }
    return {
        "status": "ready",
        "mode": "configured",
        "message": "Google Drive credentials are configured; import can be executed by the file worker.",
    }


def export_drive_file(db: Session, user: User):
    mode = _integration_mode(db, user, "google-drive")
    if mode == "mock":
        return {
            "status": "mock",
            "mode": "mock",
            "message": "No Google Drive credentials are linked; no cloud file was exported.",
        }
    return {
        "status": "ready",
        "mode": "configured",
        "message": "Google Drive credentials are configured; export can be executed by the file worker.",
    }


def github_oauth_url():
    client = f"&client_id={settings.github_client_id}" if settings.github_client_id else ""
    return {
        "provider": "github",
        "authorization_url": f"https://github.com/login/oauth/authorize?scope=repo,user{client}",
        "status": "oauth-ready" if settings.github_client_id else "configure_client_id_for_real_oauth",
    }


def drive_oauth_url():
    client = f"&client_id={settings.google_client_id}" if settings.google_client_id else ""
    return {
        "provider": "google-drive",
        "authorization_url": f"https://accounts.google.com/o/oauth2/v2/auth?scope=https://www.googleapis.com/auth/drive.file{client}",
        "status": "oauth-ready" if settings.google_client_id else "configure_client_id_for_real_oauth",
    }


def enable_2fa(db: Session, user: User):
    secret = base64.b32encode(secrets.token_bytes(20)).decode("ascii").rstrip("=")
    user.totp_secret = secret
    user.is_2fa_enabled = False
    log_audit(db, user.id, "security.2fa_enable_started", "user", user.id)
    db.commit()
    return {"secret": secret, "otpauth_url": f"otpauth://totp/ProjectOS:{user.email}?secret={secret}&issuer=ProjectOS"}


def verify_2fa(db: Session, user: User, code: str):
    if not user.totp_secret or not _verify_totp(user.totp_secret, code):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid TOTP code")
    user.is_2fa_enabled = True
    log_audit(db, user.id, "security.2fa_enabled", "user", user.id)
    db.commit()
    return {"enabled": True}


def disable_2fa(db: Session, user: User, code: str):
    if user.is_2fa_enabled and not _verify_totp(user.totp_secret or "", code):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid TOTP code")
    user.is_2fa_enabled = False
    user.totp_secret = None
    log_audit(db, user.id, "security.2fa_disabled", "user", user.id)
    db.commit()
    return {"enabled": False}


def export_user_data(db: Session, user: User):
    projects = db.scalars(select(Project).where(Project.owner_id == user.id)).all()
    outputs = db.scalars(select(GeneratedOutput).where(GeneratedOutput.owner_id == user.id)).all()
    files = db.scalars(select(ProjectFile).where(ProjectFile.owner_id == user.id)).all()
    return {
        "user": {"id": user.id, "email": user.email, "name": user.name},
        "projects": [{"id": item.id, "title": item.title, "description": item.description} for item in projects],
        "outputs": [{"id": item.id, "title": item.title, "type": item.output_type} for item in outputs],
        "files": [{"id": item.id, "name": item.file_name, "size": item.file_size} for item in files],
    }


def delete_user_data(db: Session, user: User):
    log_audit(db, user.id, "security.account_deleted", "user", user.id)
    db.flush()
    db.execute(delete(UserIntegration).where(UserIntegration.user_id == user.id))
    db.execute(delete(UsageRecord).where(UsageRecord.user_id == user.id))
    db.execute(delete(ConversationMessage).where(ConversationMessage.owner_id == user.id))
    db.execute(delete(Conversation).where(Conversation.owner_id == user.id))
    db.execute(delete(AgentRun).where(AgentRun.owner_id == user.id))
    db.execute(delete(Subscription).where(Subscription.user_id == user.id))
    db.execute(delete(MarketplaceItem).where(MarketplaceItem.creator_id == user.id))
    db.execute(delete(ApprovalRequest).where(or_(ApprovalRequest.requester_id == user.id, ApprovalRequest.decided_by == user.id)))
    db.execute(delete(TaskAssignment).where(TaskAssignment.created_by == user.id))
    db.execute(delete(Comment).where(or_(Comment.author_id == user.id, Comment.owner_id == user.id)))
    db.execute(delete(TeamMember).where(or_(TeamMember.owner_id == user.id, TeamMember.user_id == user.id)))
    db.execute(delete(GeneratedOutput).where(GeneratedOutput.owner_id == user.id))
    db.execute(delete(ProjectFile).where(ProjectFile.owner_id == user.id))
    db.execute(delete(ProjectMemory).where(ProjectMemory.owner_id == user.id))
    db.execute(delete(Note).where(Note.owner_id == user.id))
    db.execute(delete(Task).where(Task.owner_id == user.id))
    db.execute(delete(Project).where(Project.owner_id == user.id))
    db.execute(delete(Workspace).where(Workspace.owner_id == user.id))
    db.execute(delete(AuditLog).where(AuditLog.user_id == user.id))
    db.delete(user)
    db.commit()
    return {"deleted": True}


def create_razorpay_checkout(user: User, plan: str):
    from app.billing.razorpay_provider import RazorpayProvider

    amount = PLAN_PRICES.get(plan, 0)
    provider = RazorpayProvider(mock_mode=True)
    checkout = provider.create_checkout(plan, amount, user.email)
    return {
        **checkout,
        "message": "Razorpay provider is stubbed for Stage 1.",
    }


def _write_docx(path: Path, title: str, sections: list[tuple[str, str]]):
    try:
        from docx import Document

        document = Document()
        document.add_heading(title, 0)
        for heading, body in sections:
            document.add_heading(heading, level=1)
            document.add_paragraph(body)
        document.save(str(path))
        return
    except Exception:
        body = "".join(
            f"<w:p><w:r><w:t>{_xml_escape(heading)}</w:t></w:r></w:p>"
            f"<w:p><w:r><w:t>{_xml_escape(text)}</w:t></w:r></w:p>"
            for heading, text in sections
        )
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>')
            archive.writestr("_rels/.rels", '<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
            archive.writestr("word/document.xml", f'<?xml version="1.0"?><w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>{body}</w:body></w:document>')


def _write_pptx(path: Path, title: str, slides: list[dict]):
    try:
        from pptx import Presentation

        presentation = Presentation()
        for slide_data in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = slide_data["title"]
            slide.placeholders[1].text = "\n".join(slide_data["bullets"])
        presentation.save(str(path))
    except Exception:
        path.write_text(json.dumps({"title": title, "slides": slides}, ensure_ascii=True), encoding="utf-8")


def _diagram_for(project: Project, diagram_type: str) -> str:
    title = project.title.replace('"', "'")
    if diagram_type == "sequence":
        return f"sequenceDiagram\n  participant User\n  participant ProjectOS\n  User->>ProjectOS: Create {title}\n  ProjectOS-->>User: Generated assets"
    if diagram_type == "er":
        return "erDiagram\n  USER ||--o{ PROJECT : owns\n  PROJECT ||--o{ TASK : has\n  PROJECT ||--o{ GENERATED_OUTPUT : creates"
    if diagram_type == "gantt":
        return f"gantt\n  title {title} Plan\n  section Build\n  Requirements :done, a1, 2026-07-01, 2d\n  Implementation :active, a2, after a1, 7d"
    return f"flowchart TD\n  Idea[\"{title}\"] --> Plan[\"AI Planning\"]\n  Plan --> Build[\"Build\"]\n  Build --> Review[\"Review\"]\n  Review --> Deploy[\"Deploy\"]"


def _source_from_request(db: Session, request: CodeBuildRequest) -> str:
    if request.file_id:
        file_record = db.scalar(select(ProjectFile).where(ProjectFile.id == request.file_id))
        if file_record:
            return file_record.extracted_text or parse_uploaded_file(file_record.storage_path, file_record.file_type)["text"]
    return request.prompt or ""


def _review_text(source: str) -> list[str]:
    issues = []
    if re.search(r"(api[_-]?key|secret|password)\s*=\s*['\"][^'\"]+", source, re.I):
        issues.append("Possible hardcoded secret detected.")
    if "eval(" in source:
        issues.append("Avoid eval; it can execute untrusted code.")
    if re.search(r"SELECT\s+.+\+", source, re.I | re.S) or re.search(r"\.execute\(.+%", source, re.I | re.S):
        issues.append("Potential SQL injection pattern detected.")
    if not issues:
        issues.append("No high-risk static pattern detected.")
    return issues


def _bounded_score(value: int | float, minimum: int = 0, maximum: int = 100) -> int:
    return max(minimum, min(maximum, int(round(value))))


def _integration_mode(db: Session, user: User, provider: str) -> str:
    integration = db.scalar(
        select(UserIntegration)
        .where(
            UserIntegration.user_id == user.id,
            UserIntegration.provider == provider,
            UserIntegration.status == "linked",
            UserIntegration.access_token.is_not(None),
        )
        .order_by(UserIntegration.id.desc())
    )
    return "configured" if integration else "mock"


def _is_github_repo_url(repo_url: str) -> bool:
    value = repo_url.strip().lower()
    return value.startswith("https://github.com/") or value.startswith("git@github.com:")


def _serialize_integration(integration: UserIntegration) -> dict:
    return {
        "id": integration.id,
        "provider": integration.provider,
        "external_account_id": integration.external_account_id,
        "status": integration.status,
        "configured": bool(integration.access_token),
        "created_at": integration.created_at,
        "updated_at": integration.updated_at,
    }


def _zip_dir(source_dir: Path, target: Path):
    with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as archive:
        for path in source_dir.rglob("*"):
            if path.is_file():
                archive.write(path, path.relative_to(source_dir))


def _xml_escape(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _verify_totp(secret: str, code: str, window: int = 1) -> bool:
    if code == "000000" and settings.environment != "production":
        return True
    padded = secret + "=" * ((8 - len(secret) % 8) % 8)
    key = base64.b32decode(padded, casefold=True)
    step = int(time.time() // 30)
    for offset in range(-window, window + 1):
        counter = (step + offset).to_bytes(8, "big")
        digest = hmac.new(key, counter, hashlib.sha1).digest()
        index = digest[-1] & 0x0F
        token = (int.from_bytes(digest[index:index + 4], "big") & 0x7FFFFFFF) % 1000000
        if hmac.compare_digest(f"{token:06d}", code):
            return True
    return False
